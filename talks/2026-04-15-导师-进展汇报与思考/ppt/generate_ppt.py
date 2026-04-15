#!/usr/bin/env python3
"""2026-04-15 讨论 deck 生成器。"""

from __future__ import annotations

import re
import zipfile
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


WORKSPACE = Path(__file__).resolve().parent
GUIDE = WORKSPACE / "PPT_GUIDE.md"
OUTPUT = WORKSPACE / "deck.pptx"
TITLE = "2026-04-15-讨论"

FONT_SERIF = "Noto Serif CJK SC"
FONT_SANS = "Noto Sans CJK SC"
FONT_MONO = "Noto Sans Mono CJK SC"

SLIDE_IDS = [
    "s01-cover",
    "s02-agenda",
    "s03-summary",
    "s04-why-project1-now",
    "s05-four-project-map",
    "s06-project1-evidence-chain",
    "s07-baselines-overview",
    "s08-baseline-evidence",
    "s09-sources-curation",
    "s10-sources-stats",
    "s11-sources-main-types",
    "s12-sources-time-structure",
    "s13-sources-examples",
    "s14-sm-family",
    "s15-control-state-definition",
    "s16-pyfcstm-progress",
    "s17-pyfcstm-role",
    "s18-pyudbm-progress",
    "s19-infra-feedback",
    "s20-decisions-next-steps",
]


def rgb(hex_value: str) -> RGBColor:
    hex_value = hex_value.replace("#", "")
    return RGBColor.from_string(hex_value)


def clean_display(text: str) -> str:
    return text.replace("`", "")


COLORS = {
    "bg": rgb("#F5F1EA"),
    "panel": rgb("#FFFDFC"),
    "panel_alt": rgb("#F0EBE2"),
    "ink": rgb("#1F2430"),
    "muted": rgb("#68707E"),
    "navy": rgb("#203354"),
    "teal": rgb("#1D6F73"),
    "rust": rgb("#A04B39"),
    "gold": rgb("#C8912C"),
    "sage": rgb("#73876C"),
    "sand": rgb("#D8C4A3"),
    "line": rgb("#D7CFC3"),
    "white": rgb("#FFFFFF"),
    "green": rgb("#3F8F6B"),
    "amber": rgb("#D3A72A"),
    "orange": rgb("#C96B2C"),
    "blue_soft": rgb("#DDE9EE"),
    "teal_soft": rgb("#DCEDEA"),
    "rust_soft": rgb("#F2E4DE"),
    "gold_soft": rgb("#F6ECD8"),
    "sage_soft": rgb("#E3EBDD"),
}

SECTION_THEME = {
    "总体判断": COLORS["navy"],
    "project_1": COLORS["teal"],
    "文库证据": COLORS["teal"],
    "基础设施": COLORS["rust"],
    "待拍板事项": COLORS["gold"],
}


def set_paragraph_font(paragraph, *, name=FONT_SANS, size=12, bold=False, color=None, align=None):
    if align is not None:
        paragraph.alignment = align
    for run in paragraph.runs:
        run.font.name = name
        run.font.size = Pt(size)
        run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color


def set_cell_text(cell, text, *, size=11, bold=False, color=None, align=PP_ALIGN.LEFT, fill=None, font=FONT_SANS):
    cell.text = clean_display(text)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Pt(6)
    cell.margin_right = Pt(6)
    cell.margin_top = Pt(4)
    cell.margin_bottom = Pt(4)
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    paragraph = cell.text_frame.paragraphs[0]
    paragraph.alignment = align
    for run in paragraph.runs:
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color or COLORS["ink"]


def add_textbox(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    font=FONT_SANS,
    size=12,
    bold=False,
    color=None,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = valign
    frame.clear()
    lines = text.split("\n")
    for idx, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = clean_display(line)
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color or COLORS["ink"]
    return shape


def add_paragraph_box(
    slide,
    x,
    y,
    w,
    h,
    lines,
    *,
    font=FONT_SANS,
    size=12,
    color=None,
    bullet=False,
    bold_first=False,
    line_space_after=2,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.word_wrap = True
    frame.clear()
    for idx, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = clean_display(f"• {line}" if bullet else line)
        paragraph.space_after = Pt(line_space_after)
        for run in paragraph.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold_first and idx == 0
            run.font.color.rgb = color or COLORS["ink"]
    return shape


def add_panel(slide, x, y, w, h, *, fill, line=None, radius=True):
    kind = SHAPE.ROUNDED_RECTANGLE if radius else SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(1)
    return shape


def add_chip(slide, x, y, w, h, text, *, fill, color=None):
    add_panel(slide, x, y, w, h, fill=fill, line=fill, radius=True)
    add_textbox(
        slide,
        x,
        y + 0.03,
        w,
        h - 0.02,
        text,
        size=11,
        bold=True,
        color=color or COLORS["ink"],
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_card(slide, x, y, w, h, title, body_lines, *, accent, fill=None, title_size=16, body_size=11, badge=None):
    fill = fill or COLORS["panel"]
    add_panel(slide, x, y, w, h, fill=fill, line=COLORS["line"], radius=True)
    if h < 1.05:
        strip_h = min(0.1, h * 0.16)
        add_panel(slide, x, y, w, strip_h, fill=accent, line=accent, radius=True)
        if badge:
            add_chip(slide, x + 0.16, y + strip_h + 0.06, 0.46, 0.24, badge, fill=accent, color=COLORS["white"])
            title_x = x + 0.72
            title_w = w - 0.9
        else:
            title_x = x + 0.16
            title_w = w - 0.32
        add_textbox(slide, title_x, y + strip_h + 0.04, title_w, 0.22, title, font=FONT_SERIF, size=max(12, title_size - 1), bold=True, color=COLORS["navy"])
        add_paragraph_box(slide, x + 0.16, y + strip_h + 0.28, w - 0.32, max(0.22, h - strip_h - 0.32), body_lines, size=max(9.4, body_size - 0.2), color=COLORS["ink"], bullet=False)
        return
    add_panel(slide, x, y, w, 0.12, fill=accent, line=accent, radius=True)
    if badge:
        add_chip(slide, x + 0.18, y + 0.18, 0.55, 0.33, badge, fill=accent, color=COLORS["white"])
        title_x = x + 0.82
        title_w = w - 1.0
    else:
        title_x = x + 0.18
        title_w = w - 0.36
    add_textbox(slide, title_x, y + 0.18, title_w, 0.38, title, font=FONT_SERIF, size=title_size, bold=True, color=COLORS["navy"])
    add_paragraph_box(slide, x + 0.18, y + 0.63, w - 0.36, h - 0.8, body_lines, size=body_size, color=COLORS["ink"], bullet=False)


def add_context_bar(slide, x, y, w, h, cause, effect, *, accent):
    add_panel(slide, x, y, w, h, fill=COLORS["panel"], line=COLORS["line"], radius=True)
    add_chip(slide, x + 0.16, y + 0.07, 0.7, h - 0.14, "前因", fill=accent, color=COLORS["white"])
    add_textbox(slide, x + 0.96, y + 0.1, w * 0.42, h - 0.18, cause, size=10.2, color=COLORS["ink"])
    add_chip(slide, x + w * 0.54, y + 0.07, 0.7, h - 0.14, "因此", fill=COLORS["gold"], color=COLORS["ink"])
    add_textbox(slide, x + w * 0.54 + 0.82, y + 0.1, w * 0.34, h - 0.18, effect, size=10.2, color=COLORS["ink"])


def add_big_number_card(slide, x, y, w, h, value, label, detail, *, accent):
    add_panel(slide, x, y, w, h, fill=COLORS["panel"], line=COLORS["line"], radius=True)
    add_panel(slide, x, y, 0.14, h, fill=accent, line=accent, radius=True)
    add_textbox(slide, x + 0.28, y + 0.18, w - 0.45, 0.44, value, font=FONT_SERIF, size=25, bold=True, color=accent)
    add_textbox(slide, x + 0.28, y + 0.7, w - 0.45, 0.26, label, size=11, bold=True, color=COLORS["ink"])
    add_textbox(slide, x + 0.28, y + 0.98, w - 0.45, h - 1.12, detail, size=10, color=COLORS["muted"])


def add_section_background(slide, *, section, accent):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["bg"]
    add_panel(slide, 0.0, 0.0, 13.333, 0.12, fill=accent, line=accent, radius=False)
    shape = slide.shapes.add_shape(SHAPE.OVAL, Inches(11.4), Inches(-0.9), Inches(3.0), Inches(3.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["gold_soft"]
    shape.line.color.rgb = COLORS["gold_soft"]
    shape.fill.transparency = 0.22
    shape2 = slide.shapes.add_shape(SHAPE.OVAL, Inches(-0.85), Inches(5.4), Inches(2.8), Inches(2.8))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = COLORS["blue_soft"]
    shape2.line.color.rgb = COLORS["blue_soft"]
    shape2.fill.transparency = 0.12
    add_chip(slide, 11.48, 0.22, 1.35, 0.34, section, fill=accent, color=COLORS["white"])


def add_title_block(slide, title, subtitle, *, section):
    accent = SECTION_THEME[section]
    add_section_background(slide, section=section, accent=accent)
    add_textbox(slide, 0.62, 0.33, 10.8, 0.46, title, font=FONT_SERIF, size=23, bold=True, color=COLORS["navy"])
    add_textbox(slide, 0.64, 0.86, 10.7, 0.28, subtitle, size=11, color=COLORS["muted"])


def add_footer(slide, page_num, refs):
    add_panel(slide, 0.62, 6.9, 12.08, 0.01, fill=COLORS["line"], line=COLORS["line"], radius=False)
    add_textbox(slide, 0.64, 6.94, 8.8, 0.18, f"依据 {refs}", size=8.5, color=COLORS["muted"])
    add_textbox(slide, 11.95, 6.9, 0.75, 0.2, f"{page_num:02d}/20", font=FONT_MONO, size=8.5, color=COLORS["muted"], align=PP_ALIGN.RIGHT)


def add_table(slide, x, y, w, h, headers, rows, *, header_fill, col_widths=None, font_size=10.5, alignments=None):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if col_widths:
        for idx, width in enumerate(col_widths):
            table.columns[idx].width = Inches(width)
    for idx, header in enumerate(headers):
        set_cell_text(
            table.cell(0, idx),
            header,
            size=font_size,
            bold=True,
            color=COLORS["white"],
            align=PP_ALIGN.CENTER,
            fill=header_fill,
        )
    for row_idx, row in enumerate(rows, start=1):
        fill = COLORS["panel"] if row_idx % 2 else COLORS["panel_alt"]
        for col_idx, value in enumerate(row):
            align = alignments[col_idx] if alignments else (PP_ALIGN.LEFT if col_idx != len(headers) - 1 else PP_ALIGN.CENTER)
            set_cell_text(
                table.cell(row_idx, col_idx),
                value,
                size=font_size,
                color=COLORS["ink"],
                align=align,
                fill=fill,
            )
    return table


def add_chart(
    slide,
    x,
    y,
    w,
    h,
    categories,
    values,
    *,
    chart_type,
    accent,
    second_colors=None,
    max_scale=None,
):
    chart_data = CategoryChartData()
    chart_data.categories = [clean_display(cat) for cat in categories]
    chart_data.add_series("数量", values)
    chart = slide.shapes.add_chart(
        chart_type,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
        chart_data,
    ).chart
    chart.has_legend = False
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = COLORS["line"]
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.font.name = FONT_SANS
    chart.value_axis.minimum_scale = 0
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.name = FONT_SANS
    if max_scale is not None:
        chart.value_axis.maximum_scale = max_scale
    plot = chart.plots[0]
    plot.gap_width = 60
    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = accent
    series.format.line.color.rgb = accent
    if second_colors:
        for idx, color in enumerate(second_colors):
            point = series.points[idx]
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = color
            point.format.line.color.rgb = color
    return chart


def add_horizontal_segments(slide, x, y, w, h, segments, *, label_y_offset=0.18):
    total = sum(value for _, value, _, _ in segments)
    current_x = x
    for label, value, fill, text_color in segments:
        seg_w = w * value / total
        add_panel(slide, current_x, y, seg_w, h, fill=fill, line=fill, radius=False)
        add_textbox(
            slide,
            current_x,
            y + label_y_offset,
            seg_w,
            h - 0.1,
            f"{clean_display(label)}\n{value}",
            size=12,
            bold=True,
            color=text_color,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        current_x += seg_w


def add_arrow(slide, x, y, w, h, *, fill, direction="right"):
    shape_type = {
        "right": SHAPE.CHEVRON,
        "left": SHAPE.LEFT_ARROW,
        "up": SHAPE.UP_ARROW,
        "down": SHAPE.DOWN_ARROW,
    }[direction]
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    return shape


def timeline_pos(day: date, start: date, end: date, left: float, width: float) -> float:
    total_days = (end - start).days
    if total_days <= 0:
        raise ValueError("Invalid timeline range")
    return left + width * ((day - start).days / total_days)


def parse_notes_from_guide(path: Path) -> list[str]:
    text = path.read_text(encoding="utf-8").splitlines()
    notes_map: dict[str, list[str]] = {}
    current_slide = None
    current_field = None
    speaker_lines: list[str] = []
    supplement_lines: list[str] = []

    def flush():
        nonlocal speaker_lines, supplement_lines, current_slide
        if current_slide:
            chunks = []
            if speaker_lines:
                chunks.append("\n\n".join(speaker_lines))
            filtered = [item for item in supplement_lines if item and item.lower() != "none"]
            if filtered:
                chunks.append("补充：\n" + "\n".join(filtered))
            notes_map[current_slide] = ["\n\n".join(chunks).strip()]
        speaker_lines = []
        supplement_lines = []

    for line in text:
        heading = re.match(r"^###\s+(s\d{2}-[a-z0-9-]+)\s*$", line)
        if heading:
            flush()
            current_slide = heading.group(1)
            current_field = None
            continue
        if current_slide is None:
            continue
        if line.startswith("- Speaker Notes:"):
            current_field = "speaker"
            continue
        if line.startswith("- Notes Supplement:"):
            current_field = "supplement"
            continue
        if line.startswith("- ") and not line.startswith("  - "):
            current_field = None
            continue
        item = re.match(r"^  - (.+)$", line)
        if item and current_field == "speaker":
            speaker_lines.append(item.group(1).strip())
        elif item and current_field == "supplement":
            supplement_lines.append(item.group(1).strip())
    flush()

    notes = []
    for slide_id in SLIDE_IDS:
        if slide_id not in notes_map:
            raise ValueError(f"Missing notes for {slide_id} in {path}")
        notes.append(clean_display(notes_map[slide_id][0]))
    return notes


def apply_notes(output: Path, notes: list[str]) -> None:
    prs = Presentation(output)
    for slide, note in zip(prs.slides, notes, strict=True):
        slide.notes_slide.notes_text_frame.text = note
    prs.save(output)


def validate_notes(output: Path, expected_notes: list[str]) -> None:
    prs = Presentation(output)
    if len(prs.slides) != len(expected_notes):
        raise ValueError("Slide count and note count mismatch")
    for idx, (slide, expected) in enumerate(zip(prs.slides, expected_notes, strict=True), start=1):
        actual = slide.notes_slide.notes_text_frame.text.strip()
        if actual != expected.strip():
            raise ValueError(f"Notes mismatch on slide {idx}")
    with zipfile.ZipFile(output) as zf:
        notes_parts = [name for name in zf.namelist() if name.startswith("ppt/notesSlides/notesSlide")]
    if len(notes_parts) != len(expected_notes):
        raise ValueError("notesSlides part count mismatch")


def validate_slide_count(output: Path, expected: int) -> None:
    prs = Presentation(output)
    if len(prs.slides) != expected:
        raise ValueError(f"Expected {expected} slides, got {len(prs.slides)}")


def build_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["bg"]
    add_panel(slide, 0.0, 0.0, 13.333, 0.12, fill=COLORS["navy"], line=COLORS["navy"], radius=False)
    add_panel(slide, 8.55, 0.9, 3.8, 5.2, fill=COLORS["panel"], line=COLORS["line"], radius=True)
    add_panel(slide, 8.55, 0.9, 0.16, 5.2, fill=COLORS["gold"], line=COLORS["gold"], radius=True)
    add_textbox(slide, 0.75, 1.1, 6.8, 0.62, "2026-04-15-讨论", font=FONT_SERIF, size=28, bold=True, color=COLORS["navy"])
    add_textbox(slide, 0.78, 1.85, 6.5, 0.65, "当前进展、问题收束\n与下一步投稿判断", font=FONT_SERIF, size=22, bold=True, color=COLORS["ink"])
    add_textbox(slide, 0.8, 2.85, 5.8, 0.45, "这次汇报只做三件事：给出现状判断、解释为什么先收束 project_1、以及明确需要拍板的问题。", size=12, color=COLORS["muted"])
    add_chip(slide, 0.8, 4.0, 1.75, 0.38, "project_1 优先级", fill=COLORS["teal"], color=COLORS["white"])
    add_chip(slide, 2.7, 4.0, 1.6, 0.38, "pyfcstm vs pyudbm", fill=COLORS["rust"], color=COLORS["white"])
    add_chip(slide, 4.45, 4.0, 1.75, 0.38, "control-state 定义", fill=COLORS["gold"], color=COLORS["ink"])
    add_textbox(slide, 0.8, 5.2, 2.5, 0.28, "博士研究进展对齐", size=10.5, bold=True, color=COLORS["navy"])
    projects = [
        ("project_1", "建模主线 / 当前优先级最高", COLORS["teal"]),
        ("project_2", "性质与场景生成 / 接口层", COLORS["gold"]),
        ("project_3", "verification backend / 地基已起", COLORS["rust"]),
        ("project_4", "iterative repair / 依赖前置成熟", COLORS["sage"]),
    ]
    y = 1.15
    for name, desc, accent in projects:
        add_card(slide, 8.85, y, 3.1, 0.95, name, [desc], accent=accent, fill=COLORS["panel"], title_size=15, body_size=10.5)
        y += 1.1
    add_textbox(slide, 8.9, 5.62, 2.9, 0.25, "日期 2026-04-15", size=10.5, color=COLORS["muted"], align=PP_ALIGN.RIGHT)
    add_footer(slide, 1, "[1][2][7]")


def build_agenda(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(slide, "明天这次讨论，我最希望先对齐三个决策", "先拍板问题边界，再决定论文如何写和接下来 6 周怎么推", section="总体判断")
    cards = [
        ("01", "主投稿先锁定 project_1", "先把主稿打成型，不再一次性拉四个 project。", COLORS["navy"]),
        ("02", "对象先限于离散控制状态层", "先把模式、阶段、互锁、恢复、局部定时这一层讲透。", COLORS["teal"]),
        ("03", "先讲清 pyfcstm / pyudbm 分工", "一个回答目标对象，一个沉淀 timed backend。", COLORS["rust"]),
    ]
    x_positions = [0.72, 4.48, 8.24]
    for x, (badge, title, detail, accent) in zip(x_positions, cards, strict=True):
        add_card(slide, x, 1.55, 3.48, 2.0, title, [detail], accent=accent, badge=badge, title_size=15, body_size=11)
    add_textbox(slide, 0.78, 4.05, 3.2, 0.26, "后续展开顺序", size=11, bold=True, color=COLORS["navy"])
    flow = [
        ("结论先给", COLORS["navy"]),
        ("project_1 文库证据", COLORS["teal"]),
        ("基础设施与仓库进展", COLORS["rust"]),
        ("拍板事项与倒排计划", COLORS["gold"]),
    ]
    x = 0.78
    for idx, (label, fill) in enumerate(flow):
        add_panel(slide, x, 4.45, 2.85, 0.72, fill=COLORS["panel"], line=COLORS["line"], radius=True)
        add_panel(slide, x, 4.45, 0.12, 0.72, fill=fill, line=fill, radius=True)
        add_textbox(slide, x + 0.22, 4.67, 2.35, 0.2, f"{idx + 1}. {label}", size=12.5, bold=True, color=COLORS["ink"])
        if idx < len(flow) - 1:
            add_arrow(slide, x + 2.96, 4.63, 0.4, 0.34, fill=COLORS["sand"])
        x += 3.02
    add_card(
        slide,
        0.8,
        5.55,
        11.75,
        0.9,
        "一句话主线",
        ["不是再去找更多材料，而是让现有材料收敛成一篇问题边界清楚、对象清楚、基础设施落点清楚的会议论文。"],
        accent=COLORS["navy"],
        fill=COLORS["gold_soft"],
        title_size=15,
        body_size=11.5,
    )
    add_footer(slide, 2, "[1][2][7][15]")


def build_summary(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(slide, "当前最缺的不是材料，而是把问题对象收束清楚", "一页结论先给出来，后面逐页用文库与仓库证据补强", section="总体判断")
    cards = [
        ("1", "先锁定 `project_1`", "当前第一优先级是把主稿收束成型。", COLORS["navy"]),
        ("2", "真正缺的是收束", "`project_1` 当前不缺样本，缺的是问题定义。", COLORS["teal"]),
        ("3", "先做离散控制状态层", "目标对象先落到模式、阶段、互锁、恢复、局部 timer。", COLORS["gold"]),
        ("4", "pyfcstm 是研究基座", "它应被写成 control-state 基础设施。", COLORS["rust"]),
        ("5", "project_3 有地基", "pyudbm 已推进很深，仍缺核心搜索。", COLORS["sage"]),
        ("6", "反馈基础设施最关键", "LLM-based modeling 的杀手锏是持续反馈，而不只是 prompt。", COLORS["navy"]),
    ]
    positions = [
        (0.74, 1.45),
        (4.42, 1.45),
        (8.1, 1.45),
        (0.74, 3.56),
        (4.42, 3.56),
        (8.1, 3.56),
    ]
    for (x, y), (badge, title, detail, accent) in zip(positions, cards, strict=True):
        add_card(slide, x, y, 3.15, 1.72, title, [detail], accent=accent, fill=COLORS["panel"], badge=badge, title_size=15, body_size=10.5)
    add_card(
        slide,
        0.82,
        5.65,
        11.68,
        0.9,
        "Take-home",
        ["先收束对象，再拉厚实验；只要定义稳住，现有文库和仓库推进已经足以支撑一篇面向下一轮 conference 的主稿。"],
        accent=COLORS["navy"],
        fill=COLORS["teal_soft"],
        title_size=15,
        body_size=11.5,
    )
    add_footer(slide, 3, "[1][2][4][5][7][9][12][17]")


def build_why_project1_now(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(slide, "官方 2026 日程显示：A/B 主窗口已过，近端出口主要是 ESEM 与期刊", "main-track full paper 基本结束；近端只剩 ESEM、NIER 与 rolling journal", section="总体判断")
    add_context_bar(
        slide,
        0.8,
        1.25,
        11.75,
        0.42,
        "CAiSE / FM / RE / ASE / MoDELS research 等 A/B 主窗口都早于 2026-04-14。",
        "所以现在更该把 project_1 打成下一轮主稿，并把 ESEM / rolling journal 只当近端出口。",
        accent=COLORS["navy"],
    )
    add_textbox(slide, 0.78, 1.82, 4.2, 0.2, "官方 2026 A/B 会议窗口（AoE）", size=10.5, bold=True, color=COLORS["navy"])
    add_panel(slide, 0.8, 2.05, 8.35, 1.95, fill=COLORS["panel"], line=COLORS["line"], radius=True)
    timeline_left = 1.08
    timeline_width = 7.62
    timeline_start = date(2025, 11, 20)
    timeline_end = date(2026, 7, 5)
    axis_y = 2.98
    add_panel(slide, timeline_left, axis_y, timeline_width, 0.04, fill=COLORS["line"], line=COLORS["line"], radius=False)
    month_labels = [
        (date(2025, 11, 20), "2025-11"),
        (date(2025, 12, 1), "2025-12"),
        (date(2026, 1, 1), "2026-01"),
        (date(2026, 2, 1), "2026-02"),
        (date(2026, 3, 1), "2026-03"),
        (date(2026, 4, 1), "2026-04"),
        (date(2026, 5, 1), "2026-05"),
        (date(2026, 6, 1), "2026-06"),
        (date(2026, 7, 1), "2026-07"),
    ]
    for when, label in month_labels:
        x = timeline_pos(when, timeline_start, timeline_end, timeline_left, timeline_width)
        add_textbox(slide, x - 0.34, 2.22, 0.68, 0.16, label, font=FONT_MONO, size=7.2, color=COLORS["muted"], align=PP_ALIGN.CENTER)
        add_panel(slide, x - 0.008, 2.46, 0.016, 0.56, fill=COLORS["line"], line=COLORS["line"], radius=False)

    milestones = [
        (date(2025, 11, 28), 0.03, 2.56, 0.98, 0.34, "CAiSE\n11-21 / 11-28", COLORS["panel_alt"], COLORS["muted"], True),
        (date(2025, 12, 2), 0.02, 3.12, 0.94, 0.34, "FM\n11-25 / 12-02", COLORS["gold_soft"], COLORS["gold"], False),
        (date(2026, 2, 23), 0.0, 2.56, 0.96, 0.34, "RE\n02-16 / 02-23", COLORS["rust_soft"], COLORS["rust"], True),
        (date(2026, 3, 26), -0.12, 3.12, 0.88, 0.34, "ASE\n03-26", COLORS["blue_soft"], COLORS["navy"], False),
        (date(2026, 3, 27), 0.14, 2.56, 1.05, 0.34, "MoDELS\n03-20 / 03-27", COLORS["panel_alt"], COLORS["teal"], True),
        (date(2026, 5, 18), 0.0, 3.12, 1.08, 0.34, "ESEM Tech\n05-11 / 05-18", COLORS["teal_soft"], COLORS["teal"], False),
        (date(2026, 7, 1), 0.0, 2.56, 1.0, 0.34, "NIER\n06-24 / 07-01", COLORS["sage_soft"], COLORS["sage"], True),
    ]
    for when, dx, y, w, h, label, fill, accent, above in milestones:
        center_x = timeline_pos(when, timeline_start, timeline_end, timeline_left, timeline_width) + dx
        x = center_x - w / 2
        add_panel(slide, x, y, w, h, fill=fill, line=accent, radius=True)
        add_textbox(slide, x, y + 0.01, w, h - 0.01, label, size=7.4, bold=True, color=COLORS["ink"], align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if above:
            line_y = y + h
            line_h = max(0.03, axis_y - line_y)
            add_panel(slide, center_x - 0.012, line_y, 0.024, line_h, fill=accent, line=accent, radius=False)
        else:
            line_y = axis_y + 0.04
            line_h = max(0.03, y - line_y)
            add_panel(slide, center_x - 0.012, line_y, 0.024, line_h, fill=accent, line=accent, radius=False)

    current_x = timeline_pos(date(2026, 4, 14), timeline_start, timeline_end, timeline_left, timeline_width)
    add_panel(slide, current_x - 0.015, 2.46, 0.03, 0.95, fill=COLORS["navy"], line=COLORS["navy"], radius=False)
    add_textbox(slide, current_x - 0.48, 2.05, 0.96, 0.28, "当前点\n2026-04-14", font=FONT_MONO, size=8.2, color=COLORS["navy"], align=PP_ALIGN.CENTER)
    add_panel(slide, 9.35, 2.05, 3.2, 1.95, fill=COLORS["panel"], line=COLORS["line"], radius=True)
    add_textbox(slide, 9.54, 2.18, 2.82, 0.2, "仍可操作窗口", size=10.6, bold=True, color=COLORS["navy"])
    add_table(
        slide,
        9.5,
        2.42,
        2.9,
        1.08,
        ["路径", "窗口"],
        [
            ("ESEM Tech", "05-11 / 05-18"),
            ("ESEM EVR", "05-22 / 05-29"),
            ("MODELS NIER", "06-24 / 07-01"),
        ],
        header_fill=COLORS["teal"],
        col_widths=[1.35, 1.55],
        font_size=8.2,
    )
    add_textbox(slide, 9.54, 3.58, 2.8, 0.24, "ISSRE research 04-17；若 04-10 摘要未交，基本不算现实窗口。", size=8.3, color=COLORS["muted"])
    add_panel(slide, 0.8, 4.18, 7.25, 1.56, fill=COLORS["panel"], line=COLORS["line"], radius=True)
    add_textbox(slide, 0.96, 4.34, 2.2, 0.2, "当前准备度", size=10.5, bold=True, color=COLORS["navy"])
    readiness = [
        ("project_1", 0.92, COLORS["teal"], "问题、文库、baseline、pyfcstm 已够支撑主稿"),
        ("project_2", 0.56, COLORS["gold"], "接口层有意义，但当前不适合单独先发"),
        ("project_3", 0.44, COLORS["rust"], "backend 地基较深，验证原型尚未成型"),
        ("project_4", 0.24, COLORS["sage"], "依赖前几项更成熟后再起飞"),
    ]
    y = 4.66
    for label, ratio, accent, note in readiness:
        add_textbox(slide, 0.98, y, 1.0, 0.18, label, font=FONT_MONO, size=9.3, bold=True, color=COLORS["ink"])
        add_panel(slide, 2.02, y + 0.01, 2.85, 0.14, fill=COLORS["panel_alt"], line=COLORS["panel_alt"], radius=True)
        add_panel(slide, 2.02, y + 0.01, 2.85 * ratio, 0.14, fill=accent, line=accent, radius=True)
        add_textbox(slide, 5.02, y - 0.03, 2.78, 0.2, note, size=8.1, color=COLORS["muted"])
        y += 0.28
    add_panel(slide, 8.25, 4.18, 4.3, 1.56, fill=COLORS["panel"], line=COLORS["line"], radius=True)
    add_textbox(slide, 8.45, 4.34, 3.85, 0.2, "官方主页确认仍可投稿的期刊", size=10.5, bold=True, color=COLORS["navy"])
    add_paragraph_box(
        slide,
        8.45,
        4.62,
        3.82,
        0.95,
        [
            "TSE: CFP / submission",
            "SoSyM: submit your manuscript",
            "Requirements Engineering: submit your manuscript",
            "ASE Journal / EMSE: submit your manuscript",
        ],
        size=8.3,
        color=COLORS["ink"],
        bullet=False,
        line_space_after=1,
    )
    add_textbox(slide, 8.45, 5.52, 3.82, 0.16, "这些路径更适合作为近端外部输出，不应反过来牵引问题定义。", size=8.0, color=COLORS["muted"])
    add_card(
        slide,
        0.8,
        6.0,
        11.75,
        0.54,
        "页面结论",
        ["现实的 this-semester 出口主要是 ESEM 或 rolling journals；但 project_1 仍应按下一轮主稿质量来打，而不是为了赶窗口扩题。"],
        accent=COLORS["navy"],
        fill=COLORS["gold_soft"],
        title_size=14.5,
        body_size=9.8,
    )
    add_footer(slide, 4, "[1][15]")


def build_four_project_map(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(slide, "四个 project 是一条闭环，但近端主线必须先落在 project_1", "先把建模对象和目标形式主义立住，后面的 scenario、verification、repair 才有共同基座", section="总体判断")
    labels = [
        ("需求 / 描述", COLORS["panel_alt"]),
        ("建模", COLORS["teal_soft"]),
        ("性质 / 场景", COLORS["gold_soft"]),
        ("验证", COLORS["rust_soft"]),
        ("修复", COLORS["sage_soft"]),
    ]
    x = 0.85
    for idx, (label, fill) in enumerate(labels):
        add_panel(slide, x, 2.05, 2.0, 0.9, fill=fill, line=COLORS["line"], radius=True)
        add_textbox(slide, x, 2.32, 2.0, 0.2, label, size=14, bold=True, color=COLORS["ink"], align=PP_ALIGN.CENTER)
        if idx < len(labels) - 1:
            add_arrow(slide, x + 2.08, 2.31, 0.42, 0.34, fill=COLORS["sand"])
        x += 2.48
    cards = [
        (1.08, 3.55, 2.55, 1.6, "project_1", ["目标对象", "baseline", "样本与 `pyfcstm`"], COLORS["teal"]),
        (3.86, 3.55, 2.55, 1.6, "project_2", ["性质与场景", "生成接口层"], COLORS["gold"]),
        (6.64, 3.55, 2.55, 1.6, "project_3", ["profile-based verification", "timed backend"], COLORS["rust"]),
        (9.42, 3.55, 2.55, 1.6, "project_4", ["缺陷驱动", "iterative repair"], COLORS["sage"]),
    ]
    for x, y, w, h, title, lines, accent in cards:
        add_card(slide, x, y, w, h, title, lines, accent=accent, fill=COLORS["panel"], title_size=16, body_size=11)
    add_card(
        slide,
        0.85,
        5.55,
        11.8,
        0.88,
        "当前判断",
        ["博士主线仍然是生成-验证-修复闭环；只是本学期要先把建模对象与目标形式主义打稳，后续几项才有共同基座。"],
        accent=COLORS["navy"],
        fill=COLORS["gold_soft"],
        title_size=15,
        body_size=11.2,
    )
    add_footer(slide, 5, "[1][2]")


def build_project1_evidence_chain(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(slide, "project_1 的说服力来自三条文库线加上一条基础设施线", "不是平行摆设，而是共同回答“该建什么、凭什么、如何落地”", section="project_1")
    add_context_bar(
        slide,
        0.9,
        1.22,
        11.56,
        0.42,
        "project_1 不是靠单一材料就能站住，它需要比较对象、真实样本、状态机选型和可执行落地一起支持。",
        "所以这四条线必须被讲成一条证据链。",
        accent=COLORS["teal"],
    )
    add_card(slide, 0.92, 1.86, 2.55, 1.2, "baseline 比较集", ["62 篇比较集", "回答“该和谁比”"], accent=COLORS["navy"], fill=COLORS["panel"], title_size=15, body_size=10.5)
    add_card(slide, 9.85, 1.86, 2.55, 1.2, "真实样本库", ["787 篇 / 746 条正例", "回答“数据从哪里来”"], accent=COLORS["teal"], fill=COLORS["panel"], title_size=15, body_size=10.5)
    add_card(slide, 0.92, 4.16, 2.55, 1.2, "状态机家族库", ["669 + 10 条目", "回答“到底选哪类状态机”"], accent=COLORS["gold"], fill=COLORS["panel"], title_size=14.5, body_size=10.3)
    add_card(slide, 9.85, 4.16, 2.55, 1.2, "pyfcstm 基础设施", ["executable IR", "回答“如何落地”"], accent=COLORS["rust"], fill=COLORS["panel"], title_size=14.5, body_size=10.3)
    add_panel(slide, 4.15, 2.6, 5.05, 2.05, fill=COLORS["navy"], line=COLORS["navy"], radius=True)
    add_textbox(slide, 4.55, 2.9, 4.2, 0.26, "中央结论", size=12, bold=True, color=COLORS["sand"], align=PP_ALIGN.CENTER)
    add_textbox(slide, 4.4, 3.2, 4.55, 0.62, "三条文库线共同把论文对象\n压向 control-state 主问题", font=FONT_SERIF, size=18, bold=True, color=COLORS["white"], align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, 4.35, 3.9, 4.65, 0.28, "而 pyfcstm 则把这个对象变成可执行、可反馈、可持续迭代的模型空间。", size=10.3, color=COLORS["white"], align=PP_ALIGN.CENTER)
    add_arrow(slide, 3.55, 2.72, 0.45, 0.32, fill=COLORS["sand"], direction="right")
    add_arrow(slide, 8.76, 2.72, 0.45, 0.32, fill=COLORS["sand"], direction="left")
    add_arrow(slide, 3.55, 4.48, 0.45, 0.32, fill=COLORS["sand"], direction="right")
    add_arrow(slide, 8.76, 4.48, 0.45, 0.32, fill=COLORS["sand"], direction="left")
    add_footer(slide, 6, "[2][3][4][5][7]")


def build_baselines_overview(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(slide, "绿色 baseline 已经勾出了 5 条可直接讨论的方法线", "明天不该只报数量，而要把最值得讲的几篇“怎么做”讲清楚", section="文库证据")
    add_big_number_card(slide, 0.86, 1.7, 2.4, 1.16, "62", "baseline 条目", "比较对象已经够用", accent=COLORS["navy"])
    add_panel(slide, 0.86, 3.05, 2.4, 1.18, fill=COLORS["panel"], line=COLORS["line"], radius=True)
    add_textbox(slide, 1.05, 3.22, 1.95, 0.18, "状态分布", size=10.2, bold=True, color=COLORS["navy"])
    add_horizontal_segments(
        slide,
        1.02,
        3.56,
        2.08,
        0.28,
        [
            ("🟢", 14, COLORS["green"], COLORS["white"]),
            ("🟡", 19, COLORS["amber"], COLORS["ink"]),
            ("🟠", 29, COLORS["orange"], COLORS["white"]),
        ],
        label_y_offset=0.06,
    )
    add_textbox(slide, 1.04, 3.98, 2.02, 0.26, "不是继续堆数量，而是把几条最硬的方法线讲清楚。", size=9.4, color=COLORS["muted"])
    add_table(
        slide,
        3.55,
        1.62,
        8.96,
        4.92,
        ["代表论文", "方法骨架", "对 project_1 的启发"],
        [
            ("Structure/Event-Driven 2026", "多步结构分解 + 事件分解 + Hybrid 细化", "direct baseline 已经出现，问题不再是“有没有人做”"),
            ("SysML empirical 2025", "两阶段 prompting + model checking feedback repair", "真正有效的是 feedback loop，而不是一次性 prompt"),
            ("IEC 61499 2025", "iterative refinement + simulator / code generation", "一旦接上仿真与部署，方法性质就变了"),
            ("TTool AI 2024", "知识注入 + 工具链反馈 + MBSE 集成", "基础设施与领域知识会明显抬高效果"),
            ("Umple 2025", "One-shot / RAG 生成文本 DSL；zero-shot 基本失败", "小而稳的文本 DSL 仍需要示例与 schema 支撑"),
        ],
        header_fill=COLORS["navy"],
        col_widths=[2.4, 2.95, 3.61],
        font_size=9.7,
        alignments=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT],
    )
    add_panel(slide, 0.86, 5.7, 11.66, 0.58, fill=COLORS["gold_soft"], line=COLORS["gold_soft"], radius=True)
    add_textbox(slide, 1.05, 5.89, 11.2, 0.18, "页面结论：baseline 已经不是“有没有”，而是“该把哪几类方法差异和反馈路线讲清楚”。", size=10.6, bold=True, color=COLORS["navy"])
    add_footer(slide, 7, "[3]")


def build_baseline_evidence(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(slide, "baseline 的共同趋势不是更花的 prompt，而是更强的工具反馈", "几篇最值得讲的绿色条目给出的实证指向相当一致", section="文库证据")
    add_textbox(slide, 0.82, 1.38, 3.0, 0.2, "2026 direct baseline F1", size=10.5, bold=True, color=COLORS["navy"])
    add_chart(
        slide,
        0.82,
        1.62,
        3.75,
        2.45,
        ["Claude\n单提示", "GPT-4o\n单提示", "GPT-4o\nHybrid"],
        [0.7029, 0.5431, 0.6559],
        chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED,
        accent=COLORS["navy"],
        second_colors=[COLORS["navy"], COLORS["sand"], COLORS["teal"]],
        max_scale=0.85,
    )
    add_textbox(slide, 4.75, 1.38, 3.0, 0.2, "SysML empirical study 修复率 (%)", size=10.5, bold=True, color=COLORS["navy"])
    add_chart(
        slide,
        4.75,
        1.62,
        3.75,
        2.45,
        ["格式", "语法", "语义", "需求一致"],
        [94.6, 88.0, 43.1, 37.3],
        chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED,
        accent=COLORS["teal"],
        second_colors=[COLORS["teal"], COLORS["gold"], COLORS["rust"], COLORS["sage"]],
        max_scale=100,
    )
    add_card(
        slide,
        8.75,
        1.6,
        3.72,
        1.55,
        "TTool AI",
        ["状态机评分 63 vs 58", "速度快 15.2x", "块图评分 81 vs 70", "速度快 67.5x"],
        accent=COLORS["rust"],
        fill=COLORS["panel"],
        title_size=15,
        body_size=11,
    )
    add_card(
        slide,
        8.75,
        3.35,
        3.72,
        1.55,
        "IEC 61499 / workflow principles",
        ["一旦接上仿真、代码生成和 sanity checks，方法论就不再是一次 prompt。"],
        accent=COLORS["gold"],
        fill=COLORS["panel"],
        title_size=15,
        body_size=11,
    )
    add_card(
        slide,
        0.82,
        5.3,
        11.7,
        0.9,
        "Take-home",
        ["纯 prompt 可以给出骨架，但真正把质量拉上去的，是 model checking、仿真、编译检查和 workflow feedback。"],
        accent=COLORS["navy"],
        fill=COLORS["rust_soft"],
        title_size=15,
        body_size=11.5,
    )
    add_footer(slide, 8, "[12][13][14][16][17]")


def build_sources_curation(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(slide, "`sources/` 不是自然分布样本，而是面向数据集建设的治理后主集", "如果不先说明这一点，后面的统计就很容易被误读", section="文库证据")
    add_context_bar(
        slide,
        0.82,
        1.28,
        11.7,
        0.42,
        "如果不先讲治理口径，后面所有 EFSM/HSM/T0/T1 的比例都会被误读成自然分布。",
        "所以这一页先解释筛选逻辑，再看统计。",
        accent=COLORS["rust"],
    )
    add_card(
        slide,
        0.82,
        1.82,
        11.7,
        0.88,
        "Warning",
        ["后面所有 EFSM / HSM / T0 / T1 的比例，首先反映的是收录策略与数据集治理目标，其次才是领域文献现象本身。"],
        accent=COLORS["rust"],
        fill=COLORS["rust_soft"],
        title_size=15,
        body_size=11,
    )
    add_textbox(slide, 0.86, 2.78, 2.6, 0.2, "治理流程", size=10.5, bold=True, color=COLORS["navy"])
    add_card(slide, 0.86, 3.02, 2.8, 0.96, "广撒网检索", ["先把真实控制系统状态逻辑的文献面铺开。"], accent=COLORS["navy"], fill=COLORS["panel"], title_size=15, body_size=10.3)
    add_arrow(slide, 3.8, 3.32, 0.42, 0.34, fill=COLORS["sand"])
    add_card(slide, 4.28, 3.02, 2.8, 0.96, "标准化治理", ["按论文级、案例级、主类型、时间级别等口径做筛。"], accent=COLORS["teal"], fill=COLORS["panel"], title_size=15, body_size=10.3)
    add_arrow(slide, 7.22, 3.32, 0.42, 0.34, fill=COLORS["sand"])
    add_card(slide, 7.7, 3.02, 3.1, 0.96, "主集保留", ["优先保留 EFSM/HSM、T0/T1、细节高、不强趋同的案例。"], accent=COLORS["gold"], fill=COLORS["panel"], title_size=15, body_size=10.3)
    add_textbox(slide, 6.1, 4.0, 2.6, 0.2, "五套治理口径", size=10.5, bold=True, color=COLORS["navy"])
    add_table(
        slide,
        4.22,
        4.25,
        8.3,
        1.95,
        ["口径", "作用"],
        [
            ("论文级可用性", "判断单篇论文整体上能否形成可靠样本"),
            ("案例级角色", "区分核心保留、清洗后保留、降采样保留"),
            ("细节充实度", "判断是否足以支撑高质量建模与评测"),
            ("主类型", "回答“默认应把它理解成哪类状态机”"),
            ("时间级别", "回答时间语义到底只是局部 timer 还是强实时/连续耦合"),
        ],
        header_fill=COLORS["teal"],
        col_widths=[2.0, 6.3],
        font_size=10.2,
    )
    add_card(
        slide,
        0.86,
        4.25,
        3.0,
        1.95,
        "治理目标",
        [
            "优先保留 `EFSM/HSM`。",
            "优先保留 `T0/T1`。",
            "优先保留细节充实、可转为可执行样本的案例。",
            "对强趋同簇做降采样控制。",
        ],
        accent=COLORS["navy"],
        fill=COLORS["panel"],
        title_size=15,
        body_size=10.5,
    )
    add_footer(slide, 9, "[4][6]")


def build_sources_stats(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(slide, "`sources/` 的总体统计已经足以支撑数据集和问题边界判断", "统计项必须带定义列，否则数字没有解释力", section="文库证据")
    add_context_bar(
        slide,
        0.88,
        1.24,
        11.74,
        0.42,
        "样本主链已经够厚，不再是“先去找样本”的阶段。",
        "所以这页要把数量和定义一起说清楚。",
        accent=COLORS["navy"],
    )
    cards = [
        ("787", "论文总数", "正式入账论文目录", COLORS["navy"]),
        ("715", "论文级 🟢", "可直接支撑可靠样本", COLORS["green"]),
        ("746", "正例案例", "当前样本主链总量", COLORS["teal"]),
        ("685", "核心保留 `💎`", "主数据集主链", COLORS["gold"]),
    ]
    x = 0.78
    for value, label, detail, accent in cards:
        add_big_number_card(slide, x, 1.82, 2.86, 1.05, value, label, detail, accent=accent)
        x += 3.02
    add_table(
        slide,
        0.8,
        3.08,
        6.0,
        2.75,
        ["指标", "定义", "数量"],
        [
            ("论文总数", "当前 `sources/` 已正式入账的论文目录总数", "787"),
            ("论文级 `🟢`", "单篇论文整体上可直接支持可靠样本抽取", "715"),
            ("论文级 `🟡`", "单篇论文仍有价值，但需要额外整理或局部补证", "16"),
            ("论文级 `⚪`", "单篇论文最终未形成可靠样本", "56"),
        ],
        header_fill=COLORS["navy"],
        col_widths=[1.25, 3.85, 0.9],
        font_size=10.1,
    )
    add_table(
        slide,
        6.95,
        3.08,
        5.6,
        2.75,
        ["类别", "定义", "数量"],
        [
            ("正例案例总数", "正式入账、可作为控制状态样本讨论的案例条目", "746"),
            ("`💎` 核心保留", "细节和分布都足够好，适合进入主数据集主链", "685"),
            ("`🧰` 清洗后保留", "有价值，但需要进一步清洗、补证或统一口径", "20"),
            ("`🪫` 降采样保留", "本身可用，但因分布控制原因不应过量进入主集", "41"),
        ],
        header_fill=COLORS["teal"],
        col_widths=[1.35, 3.35, 0.9],
        font_size=10.1,
    )
    add_card(
        slide,
        0.82,
        6.0,
        11.7,
        0.55,
        "Take-home",
        ["主链样本已经够厚；现在真正该做的是围绕它定义论文对象、baseline 口径和可执行实验。"],
        accent=COLORS["navy"],
        fill=COLORS["gold_soft"],
        title_size=14,
        body_size=10.8,
    )
    add_footer(slide, 10, "[4][6]")


def build_sources_main_types(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(slide, "这些比例不能被读成自然分布，它们只能证明我们已经蓄了足够厚的目标样例池", "后期 sources 是按 EFSM / HSM + T0 / T1 主动强筛出来的", section="文库证据")
    add_context_bar(
        slide,
        0.84,
        1.22,
        11.62,
        0.42,
        "如果把这些比例误读成自然分布，后面的结论就会站不住。",
        "所以这一页必须明确“能说什么、不能说什么”。",
        accent=COLORS["teal"],
    )
    add_textbox(slide, 0.86, 1.86, 2.0, 0.2, "不能这样解读", size=10.8, bold=True, color=COLORS["rust"])
    add_table(
        slide,
        0.84,
        2.1,
        5.45,
        3.92,
        ["误读", "为什么不成立"],
        [
            ("“EFSM 最多”", "不能推出控制系统文献天然以 EFSM 为主；这里只能说明我们后期主动把 EFSM/HSM 样例筛得更厚。"),
            ("“T0/T1 最多”", "不能推出真实系统几乎没有强时间 / 连续问题；这里只能说明 project_1 当前主池刻意优先保留这两类。"),
            ("“429 / 157 / 127 就是总体分布”", "sources 是治理后的主样例池，不是为替整个领域做分布统计而建。"),
        ],
        header_fill=COLORS["rust"],
        col_widths=[1.48, 3.97],
        font_size=9.4,
        alignments=[PP_ALIGN.LEFT, PP_ALIGN.LEFT],
    )
    add_textbox(slide, 6.55, 1.86, 2.05, 0.2, "可以这样解读", size=10.8, bold=True, color=COLORS["navy"])
    add_table(
        slide,
        6.52,
        2.1,
        5.95,
        3.92,
        ["可以这样说", "证据", "为什么重要"],
        [
            ("目标样例池已经够厚", "FSM 127 / EFSM 429 / HSM 157", "足以支撑 project_1 的数据集与评测设计"),
            ("局部时间样例已经充足", "T0 + T1 = 719 / 746", "第一篇 paper 可以先收束到离散 control-state + local timing"),
            ("收束对象不等于样例贫瘠", "显式时钟 243 / 层次 160", "目标形式主义仍需保留 hierarchy、guard 与局部 timer"),
        ],
        header_fill=COLORS["navy"],
        col_widths=[1.86, 1.72, 2.37],
        font_size=9.35,
        alignments=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT],
    )
    add_panel(slide, 0.86, 6.12, 11.6, 0.45, fill=COLORS["teal_soft"], line=COLORS["teal_soft"], radius=True)
    add_textbox(slide, 1.04, 6.26, 11.2, 0.16, "结论：sources 的价值在于“样例池够厚”，不是“替整个领域做分布统计”。", size=10.3, bold=True, color=COLORS["navy"])
    add_footer(slide, 11, "[4]")


def build_sources_time_structure(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(slide, "这个定向样例池依然不是“简单平面 FSM 池”，它保留了层次、定时和恢复复杂度", "收敛对象不等于退化对象", section="文库证据")
    add_context_bar(
        slide,
        0.82,
        1.22,
        11.64,
        0.42,
        "如果把目标收束误读成“只做最简单 FSM”，后面 pyfcstm 的定位也会被说扁。",
        "所以这一页要强调池子里的结构复杂度。",
        accent=COLORS["navy"],
    )
    add_table(
        slide,
        0.82,
        1.82,
        11.66,
        3.98,
        ["主模式", "定义", "为什么重要", "当前证据"],
        [
            ("阶段链", "系统按步骤推进，每一步由 guard / timer 决定下一段。", "这是离散顺序控制最常见的结构骨架。", "洗衣机 PLC、电梯 PLC；T0/T1 合计 719 / 746"),
            ("联锁许可", "动作是否允许由权限、互斥、锁闭、到位条件共同决定。", "说明 guard 与变量面不能只是附属装饰。", "铁路联锁、门控控制；显式时钟 243"),
            ("模式层次", "上层模式管理下层子阶段、子任务或执行层。", "目标 DSL 至少要保留 hierarchy。", "层次样例 160；UAV supervisor 等 HSM 案例"),
            ("异常恢复", "fault -> safe fallback -> reset / recovery 构成专门链路。", "模型不能只覆盖 happy path。", "sources 中恢复链与异常保护样例长期高频出现"),
        ],
        header_fill=COLORS["navy"],
        col_widths=[1.28, 3.08, 3.0, 4.3],
        font_size=9.35,
        alignments=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT],
    )
    add_table(
        slide,
        0.82,
        5.9,
        11.66,
        0.72,
        ["结构证据", "数量", "含义"],
        [
            ("层次", "160", "离散 control-state 仍大量包含父子模式关系"),
            ("显式时钟", "243", "局部 timer / 时钟条件并不是少数样例"),
            ("T0 + T1", "719 / 746", "当前主问题可先稳定收束到离散 control-state + local timing"),
        ],
        header_fill=COLORS["teal"],
        col_widths=[1.45, 1.05, 9.16],
        font_size=9.2,
        alignments=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT],
    )
    add_panel(slide, 0.82, 6.62, 11.66, 0.32, fill=COLORS["gold_soft"], line=COLORS["gold_soft"], radius=True)
    add_textbox(slide, 1.02, 6.71, 11.2, 0.14, "结论：收敛对象 = 离散 control-state，不 = 简单平面 FSM；目标形式主义至少要能承载 EFSM + HSM + local timing。", size=9.8, bold=True, color=COLORS["navy"])
    add_footer(slide, 12, "[4][6]")


def build_sources_examples(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(slide, "代表样本显示：真实控制系统里至少并存五类不同建模难点", "同样都叫“状态机”，但它们的主难点并不相同", section="文库证据")
    add_card(
        slide,
        0.82,
        1.22,
        11.65,
        0.62,
        "Framing",
        ["这些样本不是用来“凑例子”，而是用来证明“状态机”一词背后其实罩着不同问题，因此论文对象必须主动收束。"],
        accent=COLORS["navy"],
        fill=COLORS["gold_soft"],
        title_size=14.5,
        body_size=10.6,
    )
    add_table(
        slide,
        0.82,
        2.02,
        11.65,
        4.12,
        ["场景", "标签", "关键建模难点", "对论文对象的启发"],
        [
            ("洗衣机 PLC", "EFSM / T1", "阶段链 + timer + guard", "真实工业控制里“阶段 + timer”是高频对象"),
            ("电梯 PLC", "EFSM / T1", "请求队列、方向优先、门控时长", "顺序控制与门控规则构成主问题"),
            ("铁路联锁", "Resource / T1-T2", "资源互斥、锁闭、释放", "有些“状态机”本质上更像强 guard / 资源系统"),
            ("UAV 分层任务", "HSM / T0-T1", "mission supervisor + 子层任务控制", "层次任务控制是另一条高频主线"),
            ("外骨骼步态控制", "Hybrid / T3", "连续相变量与模式切换强耦合", "确实存在第二类连续/混成问题"),
        ],
        header_fill=COLORS["navy"],
        col_widths=[1.55, 1.55, 3.45, 5.1],
        font_size=10.0,
    )
    add_card(
        slide,
        7.08,
        5.82,
        5.36,
        0.78,
        "结论",
        ["离散顺序控制、层次任务、资源互锁、连续耦合至少是四类不同问题。"],
        accent=COLORS["rust"],
        fill=COLORS["rust_soft"],
        title_size=13.5,
        body_size=9.8,
    )
    add_footer(slide, 13, "[4][6]")


def build_sm_family(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(slide, "state_machine_types 给出的启发不是“谁最多”，而是贡献点可以落在 profile / DSL / infrastructure", "现代状态机研究越来越像“选择并塑造目标对象”，而不是假定唯一标准", section="文库证据")
    add_context_bar(
        slide,
        0.84,
        1.22,
        11.62,
        0.42,
        "既然“状态机”本来就是家族，project_1 就不能假装目标对象天然唯一。",
        "所以必须主动选 control-state + infrastructure 这条分支。",
        accent=COLORS["navy"],
    )
    add_table(
        slide,
        0.84,
        1.86,
        5.72,
        4.28,
        ["主分支", "解决的复杂度", "代表对象"],
        [
            ("control-state", "模式、阶段、guard、层次组织", "FSM / EFSM / HSM"),
            ("timed", "时钟、deadline、最小/最大持续时间", "Timed Automata / clocks"),
            ("hybrid", "连续动力学与离散切换耦合", "Hybrid Automata / CPS"),
            ("interaction-resource", "资源互斥、契约组合、并发同步", "Petri / Interface / contract"),
        ],
        header_fill=COLORS["navy"],
        col_widths=[1.45, 2.2, 2.07],
        font_size=9.35,
        alignments=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT],
    )
    add_table(
        slide,
        6.76,
        1.86,
        5.7,
        4.28,
        ["贡献形态", "当前证据", "对 project_1 的意义"],
        [
            ("DSL / target profile", "59 条", "目标对象设计本身就是学术贡献位点"),
            ("标准 / 元模型 / 交换载体", "264 条", "工程上更关心可承载、可互操作、可执行"),
            ("2010+ 非模型层增量", "82.6%", "现代增量大量发生在 runtime / bridge / workbench"),
        ],
        header_fill=COLORS["teal"],
        col_widths=[1.7, 1.0, 3.0],
        font_size=9.3,
        alignments=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT],
    )
    add_panel(slide, 0.84, 6.2, 11.62, 0.34, fill=COLORS["teal_soft"], line=COLORS["teal_soft"], radius=True)
    add_textbox(slide, 1.04, 6.3, 11.2, 0.14, "结论：project_1 的贡献不一定是再发明一种“全新状态机”，也可以是主动选定并塑造一类更适合任务的 profile。", size=9.9, bold=True, color=COLORS["navy"])
    add_footer(slide, 14, "[5][7]")


def build_control_state_definition(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(slide, "我们要解的是控制系统的离散 control-state layer，而不是所有状态机", "模式组织、联锁 guard、异常恢复、事件权限和局部 timer 才是当前主问题", section="project_1")
    add_context_bar(
        slide,
        0.82,
        1.22,
        11.64,
        0.42,
        "“控制系统状态机”这个词太大，如果不拆开，第一篇 paper 就会同时背上离散、timed、hybrid 三类问题。",
        "所以这里先把当前最稳定、最可执行、最可验证的 control-state 语义块钉出来。",
        accent=COLORS["teal"],
    )
    add_table(
        slide,
        0.82,
        1.84,
        11.65,
        4.38,
        ["语义块", "工程含义", "为什么是第一篇 paper 的主问题", "pyfcstm 当前承载"],
        [
            ("模式层次", "系统有上层 mode 与下层子阶段 / 子任务。", "控制系统最常见的结构复杂度来自 hierarchy，而不是抽象图形语法。", "可直接承载"),
            ("联锁 guard", "动作是否允许取决于权限、互斥、到位、锁闭等条件。", "决定变量、guard、effect 必须成为一等建模对象。", "可直接承载"),
            ("故障与恢复链", "fault、fallback、reset、recover 构成专门路径。", "控制逻辑不能只覆盖 happy path。", "可直接承载"),
            ("事件作用域", "同一事件只在某些模式 / 边界上才有效。", "这是 hierarchy 与 control-state 语义的核心交点。", "可直接承载"),
            ("生命周期动作", "enter / during / exit 负责阶段切面与局部行为。", "比“普通状态机画图”更贴近真实控制逻辑。", "可直接承载"),
            ("局部工程定时", "延时保持、watchdog、最小停留等局部 timer。", "第一篇 paper 需要承认 local timing，但不必先吞下完整 TA / hybrid 家族。", "部分承载；强实时仍待后续映射"),
        ],
        header_fill=COLORS["navy"],
        col_widths=[1.55, 2.85, 4.15, 3.1],
        font_size=9.1,
        alignments=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT],
    )
    add_panel(slide, 0.82, 6.34, 11.65, 0.34, fill=COLORS["teal_soft"], line=COLORS["teal_soft"], radius=True)
    add_textbox(slide, 1.02, 6.44, 11.2, 0.14, "结论：第一篇 paper 先解这六类 control-state 语义；连续 / 混成控制的问题承认其重要性，但放到后续延展更稳。", size=9.8, bold=True, color=COLORS["navy"])
    add_footer(slide, 15, "[4][5][6][7]")


def build_pyfcstm_progress(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(slide, "在 STM 语境下，pyfcstm 更像一个收窄后的 control-state DSL，而不是大而全状态机工具", "最该重点比较的是 Umple / UmpleRun 这条文本状态机生态，而不是泛泛说“它像状态图”", section="基础设施")
    add_panel(slide, 0.84, 1.32, 11.64, 0.52, fill=COLORS["panel"], line=COLORS["line"], radius=True)
    add_textbox(slide, 1.02, 1.48, 11.2, 0.16, "一句话定位：sequential HSM 骨架 + EFSM 数据面 + 确定执行语义的 executable control-state DSL。", size=10.8, bold=True, color=COLORS["navy"])
    add_table(
        slide,
        0.84,
        2.02,
        11.64,
        4.42,
        ["近邻对象", "它在 STM 文库里做什么", "与 pyfcstm 的相似点", "关键差异"],
        [
            ("HSM / Statecharts", "层次控制与复合状态的家族血缘", "都把 hierarchy 当结构骨架", "pyfcstm 主动收窄语义面，不追求完整并发 / history / 广播语义"),
            ("Umple", "文本化 UML / 复合状态机 DSL 与代码生成", "都强调文本状态机、层次、guard / action 与可执行工件", "Umple 更偏 UML 文本承载与 Java/C++ 代码生成；pyfcstm 更强调 control-state profile 与 formal core"),
            ("UmpleRun", "Umple 生态里的 execution-scenario 动态验证", "都说明文本状态机可以直接接反馈闭环", "UmpleRun 依赖 Umple -> Java/JAR -> scenario；pyfcstm 直接把 parser、runtime、symbolic core 放在同一内核里"),
            ("SCXML", "标准化 executable state machine / interchange 载体", "都属于可执行层次状态机文本承载", "SCXML 语义面更宽，含 queue、parallel、history、invoke/send；pyfcstm 更窄、更偏 control-state 闭核"),
            ("Sismic", "statechart runtime / testing 路线", "都强调 executable state machine 与测试反馈", "Sismic 更像 Python runtime；pyfcstm 更强调形式化核心与外部副作用隔离"),
            ("UPPAAL", "timed automata / verification backend", "都和后续验证闭环强相关", "UPPAAL 属于 clocks / invariants / symbolic reachability 家族；pyfcstm 当前仍是 control-state DSL 本体"),
        ],
        header_fill=COLORS["rust"],
        col_widths=[1.45, 2.35, 2.9, 4.94],
        font_size=8.75,
        alignments=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT],
    )
    add_panel(slide, 0.84, 6.54, 11.64, 0.16, fill=COLORS["gold_soft"], line=COLORS["gold_soft"], radius=True)
    add_textbox(slide, 1.02, 6.57, 11.2, 0.12, "结论：pyfcstm 的辨识度来自“窄 profile + 闭 formal core”，而不是“大而全兼容”。", size=9.8, bold=True, color=COLORS["navy"])
    add_footer(slide, 16, "[5][7][8]")


def build_pyfcstm_role(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(slide, "pyfcstm 的研究价值，在于把目标形式主义、执行语义和闭环接口一起做出来", "它不是附属工具，而是 project_1 的核心研究产出之一", section="基础设施")
    add_panel(slide, 0.84, 1.32, 11.62, 0.5, fill=COLORS["panel"], line=COLORS["line"], radius=True)
    add_textbox(slide, 1.02, 1.47, 11.2, 0.16, "它不是“又做了一个状态机工具”，而是在回答：应该生成什么对象、为什么它能立刻运行、以及为什么它能继续接向验证与修复闭环。", size=10.2, bold=True, color=COLORS["navy"])
    add_table(
        slide,
        0.84,
        2.02,
        5.6,
        4.16,
        ["贡献点", "真正回答的问题"],
        [
            ("target profile", "LLM 应该生成什么对象，而不是任意 UML / SCXML 子集"),
            ("executable semantics", "为什么生成结果一落地就能跑起来，而不是只是一张图"),
            ("formal core boundary", "哪些语义留在模型核心里，哪些必须外挂给 abstract action"),
            ("executable model output", "输出对象为何从状态图升级成可执行形式模型"),
            ("unified analysis-ready substrate", "为什么它能继续接 parser / runtime / verification hooks"),
        ],
        header_fill=COLORS["navy"],
        col_widths=[1.65, 3.95],
        font_size=9.2,
        alignments=[PP_ALIGN.LEFT, PP_ALIGN.LEFT],
    )
    add_table(
        slide,
        6.86,
        2.02,
        5.6,
        4.16,
        ["当前能力层", "已有落地", "为什么重要"],
        [
            ("parser", "DSL 解析 + 结构校验", "保证目标对象不是松散文本"),
            ("runtime", "cycle、stable-boundary、rollback", "为仿真与反馈提供确定执行边界"),
            ("symbolic expr", "表达式到 Z3、symbolic execution", "为后续 analysis-ready 路线预埋接口"),
            ("codegen", "Python / C 模板与生成路线", "让模型能继续落向实现侧"),
            ("tooling", "PlantUML、文档、教程、编辑器支持", "让它成为可维护对象而不是 demo"),
        ],
        header_fill=COLORS["rust"],
        col_widths=[1.15, 1.75, 2.7],
        font_size=8.9,
        alignments=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT],
    )
    add_panel(slide, 0.84, 6.28, 11.62, 0.34, fill=COLORS["gold_soft"], line=COLORS["gold_soft"], radius=True)
    add_textbox(slide, 1.04, 6.39, 11.2, 0.12, "结论：pyfcstm 是 project_1 对“应该生成什么、怎样立刻可用、如何继续进闭环”的研究性回答。", size=9.8, bold=True, color=COLORS["navy"])
    add_footer(slide, 17, "[7][8]")


def build_pyudbm_progress(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(slide, "`project_3` 的实质推进主要在 `pyudbm`，但核心搜索仍缺", "backend 不是空白，只是还没长成第一版 profile-guided verifier", section="基础设施")
    add_context_bar(
        slide,
        0.82,
        1.22,
        11.64,
        0.42,
        "project_3 很容易被看成“目录还空，所以几乎没动”。",
        "所以这页要把“地基已做很多”和“完整 verifier 仍未成型”分开说。",
        accent=COLORS["rust"],
    )
    add_big_number_card(slide, 0.82, 1.76, 2.3, 1.02, "a8d0649", "main HEAD", "本地 clone 于 `~/oo-projects/pyudbm`", accent=COLORS["rust"])
    stack_items = [
        ("符号核", "UDBM API 与 UCDD 路线已立住", COLORS["navy"]),
        ("模型前端", "UTAP 绑定与 query 接口已接通", COLORS["teal"]),
        ("query + corpus", "roundtrip 与 178 个官方样本已成体系", COLORS["gold"]),
        ("文献与路线", "TA / UPPAAL 阅读地图持续维护", COLORS["sage"]),
    ]
    y = 2.95
    for idx, (title, detail, accent) in enumerate(stack_items):
        width = 3.1 + idx * 0.34
        add_panel(slide, 0.98, y, width, 0.52, fill=COLORS["panel"], line=accent, radius=True)
        add_textbox(slide, 1.18, y + 0.15, width - 0.3, 0.22, f"{title}  |  {detail}", size=10.2, bold=True, color=COLORS["ink"])
        y += 0.64
    add_table(
        slide,
        6.15,
        1.88,
        6.3,
        4.18,
        ["层", "已有", "仍缺 / 影响"],
        [
            ("基础符号层", "已有", "不是主要瓶颈"),
            ("模型 / query 接口", "已有", "官方样本和 roundtrip 已成体系"),
            ("symbolic reachability", "仍缺", "无法形成真正 verifier"),
            ("A[] / E<> 求值", "仍缺", "性质检查闭环还接不上"),
            ("witness / counterexample", "仍缺", "反馈链条不完整"),
            ("verifyta 核心搜索", "仍缺", "这才是 project_3 真正未过的门槛"),
        ],
        header_fill=COLORS["rust"],
        col_widths=[1.8, 1.0, 3.5],
        font_size=10.0,
    )
    add_card(
        slide,
        0.84,
        6.02,
        11.62,
        0.74,
        "页面结论",
        ["`project_3` 现在最合理的定位，是继续让 `pyudbm` 沉淀 timed backend，而不是抢在 `project_1` 之前承担主投稿。"],
        accent=COLORS["rust"],
        fill=COLORS["rust_soft"],
        title_size=13.5,
        body_size=10.2,
    )
    add_footer(slide, 18, "[9][10][11]")


def build_infra_feedback(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(slide, "我越来越相信基础设施反馈才是 LLM-based modeling 的真正杀手锏", "parser、runtime、simulation、model checking 和 regression traceability 会把问题性质彻底改变", section="基础设施")
    add_context_bar(
        slide,
        0.82,
        1.22,
        11.64,
        0.42,
        "单次生成很快会碰到 action、语义一致性和可执行性瓶颈。",
        "所以真正把质量拉高的，不是 prompt 花样，而是反馈基础设施。",
        accent=COLORS["navy"],
    )
    add_panel(slide, 3.55, 2.42, 6.15, 2.18, fill=COLORS["panel"], line=COLORS["line"], radius=True)
    loop_boxes = [
        (4.0, 2.78, 1.1, 0.56, "LLM", COLORS["gold_soft"]),
        (5.38, 2.78, 1.55, 0.56, "control-state\nDSL", COLORS["teal_soft"]),
        (7.28, 2.78, 1.55, 0.56, "parser /\nruntime", COLORS["blue_soft"]),
        (7.28, 3.48, 1.75, 0.64, "simulator /\nchecker /\ntests", COLORS["rust_soft"]),
        (5.15, 3.48, 1.75, 0.64, "structured\nfeedback", COLORS["sage_soft"]),
    ]
    for x, y, w, h, label, fill in loop_boxes:
        add_panel(slide, x, y, w, h, fill=fill, line=COLORS["line"], radius=True)
        add_textbox(slide, x, y + 0.12, w, h - 0.08, label, size=10.3, bold=True, color=COLORS["ink"], align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_arrow(slide, 5.08, 2.95, 0.24, 0.22, fill=COLORS["sand"], direction="right")
    add_arrow(slide, 6.98, 2.95, 0.24, 0.22, fill=COLORS["sand"], direction="right")
    add_arrow(slide, 7.88, 3.28, 0.22, 0.22, fill=COLORS["sand"], direction="down")
    add_arrow(slide, 6.88, 3.74, 0.26, 0.22, fill=COLORS["sand"], direction="left")
    add_arrow(slide, 4.38, 3.12, 0.22, 0.36, fill=COLORS["sand"], direction="up")
    evidence_cards = [
        (0.82, 1.82, 2.3, 1.12, "2026 direct baseline", ["0.7029 vs 0.5431", "Hybrid 流程能拉回弱模型"], COLORS["navy"]),
        (10.2, 1.82, 2.25, 1.18, "SysML empirical", ["94.6 / 88.0 / 43.1 / 37.3", "规则反馈有效，但还不够"], COLORS["teal"]),
        (0.82, 4.7, 2.3, 1.12, "TTool AI", ["63 vs 58，15.2x", "81 vs 70，67.5x"], COLORS["rust"]),
        (10.2, 4.62, 2.25, 1.2, "workflow principles", ["可信 GenAI 依赖 decomposition + checks + traceability"], COLORS["gold"]),
    ]
    for x, y, w, h, title, lines, accent in evidence_cards:
        add_card(slide, x, y, w, h, title, lines, accent=accent, fill=COLORS["panel"], title_size=13.5, body_size=10)
    add_card(
        slide,
        3.15,
        5.42,
        7.1,
        0.88,
        "Take-home",
        ["与其继续拼 prompt，不如把模型输出放进能持续打分、回传反例、约束格式和校验语义的环境里。"],
        accent=COLORS["navy"],
        fill=COLORS["gold_soft"],
        title_size=14.5,
        body_size=11,
    )
    add_footer(slide, 19, "[12][13][14][16][17]")


def build_decisions_next_steps(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_block(slide, "如果明天拍板这五件事，后续 6 周的推进路径会更稳", "收束问题对象、拉厚实验、准备 ESEM / rolling journal 后手", section="待拍板事项")
    add_context_bar(
        slide,
        0.84,
        1.22,
        11.62,
        0.42,
        "如果前面几项关键判断不能拍板，后面 6 周就会继续分散推进。",
        "所以这一页把倒排计划和待确认问题放在一起。",
        accent=COLORS["gold"],
    )
    add_textbox(slide, 0.84, 1.76, 2.0, 0.2, "6 周倒排", size=10.5, bold=True, color=COLORS["navy"])
    phases = [
        ("第 1-2 周", "收束问题定义与论文主张", COLORS["teal"]),
        ("第 3-4 周", "固化样本与 baseline 实验", COLORS["gold"]),
        ("第 5-6 周", "完成初稿并锁定 ESEM / journal 后手", COLORS["rust"]),
    ]
    x = 0.86
    for idx, (label, detail, accent) in enumerate(phases):
        add_card(slide, x, 2.02, 3.55, 1.0, label, [detail], accent=accent, fill=COLORS["panel"], title_size=14.5, body_size=10.4)
        if idx < len(phases) - 1:
            add_arrow(slide, x + 3.65, 2.3, 0.38, 0.32, fill=COLORS["sand"])
        x += 3.95
    decisions = [
        "主投稿是否锁定 project_1",
        "对象是否先限于离散控制状态层",
        "pyfcstm 是否写成目标形式主义",
        "pyudbm 是否继续沉淀 backend 地基",
        "主稿是否按 feedback infrastructure 组织，并以 ESEM / journal 作近端出口",
    ]
    x_positions = [0.86, 4.05, 7.24, 0.86, 4.05]
    y_positions = [3.4, 3.4, 3.4, 5.0, 5.0]
    accents = [COLORS["navy"], COLORS["teal"], COLORS["gold"], COLORS["rust"], COLORS["sage"]]
    for idx, (x, y, text, accent) in enumerate(zip(x_positions, y_positions, decisions, accents, strict=True), start=1):
        add_card(slide, x, y, 2.95, 1.2, f"决策 {idx}", [text], accent=accent, fill=COLORS["panel"], title_size=14, body_size=10.0)
    add_card(
        slide,
        10.22,
        4.1,
        2.25,
        2.1,
        "收束判断",
        ["只要 3-4 项能拍板，后面 6 周就能按“对象收束 -> 实验拉厚 -> 初稿成型 -> ESEM / journal 选择”推进。"],
        accent=COLORS["navy"],
        fill=COLORS["gold_soft"],
        title_size=15,
        body_size=10.5,
    )
    add_footer(slide, 20, "[1][7][15]")


def build_presentation() -> Path:
    notes = parse_notes_from_guide(GUIDE)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = TITLE
    prs.core_properties.subject = "讨论 deck"
    prs.core_properties.author = "Codex"

    build_cover(prs)
    build_agenda(prs)
    build_summary(prs)
    build_why_project1_now(prs)
    build_four_project_map(prs)
    build_project1_evidence_chain(prs)
    build_baselines_overview(prs)
    build_baseline_evidence(prs)
    build_sources_curation(prs)
    build_sources_stats(prs)
    build_sources_main_types(prs)
    build_sources_time_structure(prs)
    build_sources_examples(prs)
    build_sm_family(prs)
    build_control_state_definition(prs)
    build_pyfcstm_progress(prs)
    build_pyfcstm_role(prs)
    build_pyudbm_progress(prs)
    build_infra_feedback(prs)
    build_decisions_next_steps(prs)

    prs.save(OUTPUT)
    validate_slide_count(OUTPUT, 20)
    apply_notes(OUTPUT, notes)
    validate_notes(OUTPUT, notes)
    return OUTPUT


if __name__ == "__main__":
    result = build_presentation()
    print(result)
