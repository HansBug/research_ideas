#!/usr/bin/env python3
"""
讨论工作区初始化工具

用于在 talks/ 下初始化单次讨论子目录，默认同时创建：
- prep/ 准备材料
- ppt/ 讲解材料与生成器
- raw.md 会后原始碎片
- minutes.md 最终纪要
- todo.md 后续动作
"""

import argparse
import sys
from pathlib import Path
from textwrap import dedent


def derive_title(slug: str, provided_title: str | None) -> str:
    if provided_title:
        return provided_title.strip()
    return slug.replace("-", " ").strip()


def write_file(path: Path, content: str, force: bool) -> bool:
    if path.exists() and not force:
        return False
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content, encoding="utf-8")
    return True


def prep_notes_template(title: str) -> str:
    return dedent(
        f"""\
        # {title} 准备笔记

        ## 1. 这次讨论想解决什么

        - 

        ## 2. 当前我的判断

        - 

        ## 3. 想向对方确认的问题

        - 

        ## 4. 希望带去的结论或材料

        - 
        """
    )


def prep_materials_template() -> str:
    return dedent(
        """\
        # 准备材料清单

        ## 1. 需要提前回看的路径

        - 

        ## 2. 需要带去的图表、论文或数据

        - 

        ## 3. 备注

        - 
        """
    )


def raw_template(title: str) -> str:
    return dedent(
        f"""\
        # {title} 原始记录

        > 这里保留会后第一时间写下的原始片段，不要求成句，不要提前润色。

        ## 1. 记忆片段

        - 

        ## 2. 不确定但值得保留的点

        - 
        """
    )


def minutes_template(title: str) -> str:
    return dedent(
        f"""\
        # {title} 讨论纪要

        ## 1. 讨论背景

        - 

        ## 2. 已确认结论

        - 

        ## 3. 展开说明

        - 

        ## 4. 待确认点

        - 

        ## 5. 后续动作

        - 
        """
    )


def todo_template() -> str:
    return dedent(
        """\
        # 后续动作

        - [ ] 
        """
    )


def ppt_guide_template(title: str) -> str:
    return dedent(
        f"""\
        # {title} PPT 指南

        ## 1. 目标

        - 这份 deck 要帮助谁理解什么。

        ## 2. 受众

        - 

        ## 3. 讲述主线

        - 

        ## 4. 幻灯片规划

        ### slide-01-cover

        - 主要信息：
        - 屏幕文本：
        - 讲者备注：
        - 验收标准：

        ## 5. 需要准备的图表或素材

        - 

        ## 6. 约束

        - 使用 `generate_ppt.py` 作为唯一生成入口。
        - 讲稿、结构和备注优先回写本文件，再重建 `deck.pptx`。
        """
    )


def ppt_generator_template(title: str) -> str:
    safe_title = title.replace("\\", "\\\\").replace('"', '\\"')
    return dedent(
        f"""\
        #!/usr/bin/env python3
        \"\"\"{title} PPT 生成器。\"\"\"

        from pathlib import Path

        from pptx import Presentation


        WORKSPACE = Path(__file__).resolve().parent
        OUTPUT = WORKSPACE / "deck.pptx"
        TITLE = "{safe_title}"
        SUBTITLE = "TODO: 根据 PPT_GUIDE.md 补全内容后重新生成"


        def build_cover(prs: Presentation) -> None:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = TITLE
            slide.placeholders[1].text = SUBTITLE


        def build_presentation() -> Path:
            prs = Presentation()
            prs.core_properties.title = TITLE
            build_cover(prs)
            prs.save(str(OUTPUT))
            return OUTPUT


        if __name__ == "__main__":
            print(build_presentation())
        """
    )


def review_notes_template() -> str:
    return dedent(
        """\
        # Review Notes

        | Slide ID | Issue | Route To | Status |
        |---|---|---|---|
        """
    )


def review_commands_template(slug: str) -> str:
    return dedent(
        f"""\
        # 渲染与检查命令

        ## 1. 使用仓库自己的 Python 环境

        ```bash
        source venv/bin/activate
        pip install -r requirements.txt
        python talks/{slug}/ppt/generate_ppt.py
        ```

        ## 2. 渲染 review 产物

        ```bash
        python ~/.codex/skills/deck-workflow/scripts/render_review.py \\
          talks/{slug}/ppt/deck.pptx \\
          --output-dir talks/{slug}/ppt/rendered
        ```
        """
    )


def create_workspace(workspace: Path, title: str, force: bool) -> list[Path]:
    created: list[Path] = []
    for directory in [
        workspace / "prep",
        workspace / "ppt" / "assets",
        workspace / "ppt" / "rendered",
        workspace / "ppt" / "review",
    ]:
        directory.mkdir(parents=True, exist_ok=True)

    templates = {
        workspace / "prep" / "notes.md": prep_notes_template(title),
        workspace / "prep" / "materials.md": prep_materials_template(),
        workspace / "raw.md": raw_template(title),
        workspace / "minutes.md": minutes_template(title),
        workspace / "todo.md": todo_template(),
        workspace / "ppt" / "PPT_GUIDE.md": ppt_guide_template(title),
        workspace / "ppt" / "generate_ppt.py": ppt_generator_template(title),
        workspace / "ppt" / "review" / "notes.md": review_notes_template(),
        workspace / "ppt" / "review" / "commands.md": review_commands_template(
            workspace.name
        ),
        workspace / "ppt" / "assets" / ".gitkeep": "",
        workspace / "ppt" / "rendered" / ".gitkeep": "",
    }

    for path, content in templates.items():
        if write_file(path, content, force):
            created.append(path)

    return created


def main() -> None:
    parser = argparse.ArgumentParser(
        description="初始化 talks/ 下的单次讨论工作区"
    )
    parser.add_argument(
        "slug",
        help="讨论目录名，推荐 yyyy-mm-dd-对象-主题；同日多次讨论再追加 hh-mm",
    )
    parser.add_argument(
        "--root",
        default="talks",
        help="讨论根目录，默认是 talks",
    )
    parser.add_argument(
        "--title",
        default=None,
        help="讨论标题，默认由 slug 推导",
    )
    parser.add_argument(
        "--force",
        action="store_true",
        help="覆盖已存在的模板文件",
    )

    args = parser.parse_args()

    slug = args.slug.strip().strip("/")
    if not slug:
        print("错误: slug 不能为空", file=sys.stderr)
        sys.exit(1)
    if "/" in slug:
        print("错误: slug 不应包含路径分隔符，请只传目录名", file=sys.stderr)
        sys.exit(1)

    workspace = Path(args.root) / slug
    title = derive_title(slug, args.title)
    created = create_workspace(workspace, title, args.force)

    print(f"已初始化讨论工作区: {workspace}")
    if created:
        print("新建文件:")
        for path in created:
            print(f"- {path}")
    else:
        print("未新建文件；如需覆盖模板，请加 --force")

    print("\n建议下一步:")
    print(f"1. 先完善 {workspace / 'prep' / 'notes.md'}")
    print(f"2. 再完善 {workspace / 'ppt' / 'PPT_GUIDE.md'}")
    print(f"3. 运行 python {workspace / 'ppt' / 'generate_ppt.py'}")
    print(f"4. 讨论后把原始片段写入 {workspace / 'raw.md'}")
    print(f"5. 最后整理 {workspace / 'minutes.md'}")


if __name__ == "__main__":
    main()
